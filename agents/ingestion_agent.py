# agents/ingestion_agent.py

import pdfplumber
import pandas as pd
from docx import Document
from pptx import Presentation
import markdown

import os

class IngestionAgent:
    def __init__(self):
        pass

    def extract_text_from_file(self, file_path: str) -> str:
        """
        Parses a file based on its extension and extracts raw text.
        """
        ext = os.path.splitext(file_path)[-1].lower()
        content = ""

        try:
            if ext == ".pdf":
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        content += page.extract_text() + "\n"

            elif ext == ".docx":
                doc = Document(file_path)
                for para in doc.paragraphs:
                    content += para.text + "\n"

            elif ext == ".pptx":
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"

            elif ext == ".csv":
                df = pd.read_csv(file_path)
                content += df.to_string(index=False)

            elif ext in [".txt", ".md"]:
                with open(file_path, "r", encoding="utf-8") as f:
                    content += f.read()

            else:
                content = "[Unsupported file format]"

        except Exception as e:
            content = f"[Error reading {file_path}: {str(e)}]"

        return content
